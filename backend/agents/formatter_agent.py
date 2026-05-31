import json
import subprocess
import sys
import os
from datetime import date

COLOUR_HEX = {"GREEN": "1DB954", "AMBER": "F5A623", "RED": "E53935"}
COLOUR_EMOJI = {"GREEN": "🟢", "AMBER": "🟡", "RED": "🔴"}


def format_slack(deals: list[dict]) -> dict:
    
    blocks = [
        {
            "type": "header",
            "text": {"type": "plain_text", "text": "📊 Pipeline Scorecard", "emoji": True},
        },
        {"type": "divider"},
    ]

    for d in deals:
        emoji = COLOUR_EMOJI[d["colour"]]
        blocks += [
            {
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": f"*{emoji} {d['deal_name']}*",
                },
                "fields": [
                    {"type": "mrkdwn", "text": f"*Value:*\n${d['value']:,.0f}"},
                    {"type": "mrkdwn", "text": f"*Close:*\n{d['close_date']}"},
                    {"type": "mrkdwn", "text": f"*Score:*\n{d['score']}%"},
                    {"type": "mrkdwn", "text": f"*Status:*\n{d['colour']}"},
                ],
            },
            {
                "type": "context",
                "elements": [{"type": "mrkdwn", "text": d["summary"]}],
            },
            {"type": "divider"},
        ]

    return {"blocks": blocks}



def format_docx(deals: list[dict], output_path: str) -> str:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    COLOUR_RGB = {
        "GREEN": RGBColor(0x1D, 0xB9, 0x54),
        "AMBER": RGBColor(0xF5, 0xA6, 0x23),
        "RED":   RGBColor(0xE5, 0x39, 0x35),
    }

    doc = Document()
    doc.add_heading("Sales Pipeline Scorecard", 0)
    doc.add_paragraph(f"Generated: {date.today()}")

    # Summary table
    table = doc.add_table(rows=1, cols=5)
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for i, h in enumerate(["Deal", "Value", "Close", "Score %", "Status"]):
        hdr[i].text = h

    for d in deals:
        row = table.add_row().cells
        row[0].text = d["deal_name"]
        row[1].text = f"${d['value']:,.0f}"
        row[2].text = d["close_date"]
        row[3].text = f"{d['score']}%"
        row[4].text = d["colour"]

    doc.add_paragraph("")

    # Per-deal detail
    for d in deals:
        h = doc.add_heading(d["deal_name"], level=2)
        h.runs[0].font.color.rgb = COLOUR_RGB[d["colour"]]
        doc.add_paragraph(
            f"Value: ${d['value']:,.0f}   Close: {d['close_date']}   Score: {d['score']}%"
        )
        doc.add_paragraph(d["summary"])
        doc.add_paragraph("")

    doc.save(output_path)
    return output_path



def format_salesforce_json(deals: list[dict]) -> list[dict]:
    """
    Builds a list of Salesforce Opportunity update payloads.

    Maps our internal fields to standard Salesforce Opportunity API fields:
      deal_name  → Name
      value      → Amount
      close_date → CloseDate
      colour     → RAG_Status__c  (custom field)
      score      → RAG_Score__c   (custom field)
      summary    → Description

    To push to Salesforce, iterate the returned list and call:
        PATCH /services/data/v58.0/sobjects/Opportunity/{Id}
    with each payload (after resolving the opportunity Id by Name).
    """
    payloads = []
    for d in deals:
        payloads.append({
            "Name":          d["deal_name"],
            "Amount":        d["value"],
            "CloseDate":     d["close_date"],        
            "Description":   d["summary"],

            "RAG_Status__c": d["colour"],            
            "RAG_Score__c":  d["score"],             

            "_meta": {
                "source":    "formatter_agent",
                "generated": str(date.today()),
            }
        })
    return payloads



def format_pptx(deals: list[dict], output_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    COLOUR_RGB = {
        "GREEN": RGBColor(0x1D, 0xB9, 0x54),
        "AMBER": RGBColor(0xF5, 0xA6, 0x23),
        "RED":   RGBColor(0xE5, 0x39, 0x35),
    }

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    def add_textbox(slide, text, x, y, w, h, size=14, bold=False, colour=None):
        txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if colour:
            run.font.color.rgb = colour
        return txb

    # ── Slide 1: summary table ────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_textbox(s1, "Pipeline Scorecard", 0.4, 0.2, 12, 0.8, size=36, bold=True)
    add_textbox(s1, str(date.today()), 0.4, 0.95, 12, 0.35, size=12,
                colour=RGBColor(0x88, 0x88, 0x99))

    headers = ["Deal", "Value", "Close", "Score", "Status"]
    col_x   = [0.4, 4.3, 6.4, 8.5, 10.6]
    col_w   = [3.7, 1.9, 1.9, 1.9, 2.5]

    for i, h in enumerate(headers):
        add_textbox(s1, h, col_x[i], 1.5, col_w[i], 0.4, size=12, bold=True)

    for row_i, d in enumerate(deals):
        y = 2.0 + row_i * 0.55
        vals = [d["deal_name"], f"${d['value']:,.0f}", d["close_date"],
                f"{d['score']}%", d["colour"]]
        for i, v in enumerate(vals):
            colour = COLOUR_RGB[d["colour"]] if i == 4 else None
            add_textbox(s1, v, col_x[i], y, col_w[i], 0.45, size=12,
                        bold=(i == 4), colour=colour)

    # ── Slides 2-N: one per deal ───────────────────────────────────────────────
    for d in deals:
        sl = prs.slides.add_slide(blank)
        rgb = COLOUR_RGB[d["colour"]]

        add_textbox(sl, d["deal_name"], 0.4, 0.2, 10, 0.75,
                    size=32, bold=True)
        add_textbox(sl, d["colour"],   11.0, 0.25, 2.0, 0.55,
                    size=16, bold=True, colour=rgb)

        metrics = [("VALUE", f"${d['value']:,.0f}"),
                   ("CLOSE DATE", d["close_date"]),
                   ("SCORE", f"{d['score']}%")]
        for i, (label, val) in enumerate(metrics):
            x = 0.4 + i * 4.2
            add_textbox(sl, label, x, 1.2, 3.8, 0.3, size=10,
                        colour=RGBColor(0x88, 0x88, 0x99))
            add_textbox(sl, val,   x, 1.5, 3.8, 0.5, size=22, bold=True)

        add_textbox(sl, "DEAL SUMMARY", 0.4, 2.3, 12.5, 0.35,
                    size=11, bold=True, colour=rgb)
        add_textbox(sl, d["summary"],   0.4, 2.75, 12.5, 2.5, size=13)

    prs.save(output_path)
    return output_path



def run_formatter_agent(summarised_deals: list[dict], output_dir: str =os.getcwd()) -> dict:
    
    print("\n── Formatter Agent ────────────────────────────────────")

    print("  [1/4] Building Slack Block Kit payload...")
    slack_payload = format_slack(summarised_deals)
    print(f"        → {len(slack_payload['blocks'])} blocks generated")

    print("  [2/4] Building Salesforce JSON payloads...")
    sf_payloads = format_salesforce_json(summarised_deals)
    print(f"        → {len(sf_payloads)} opportunity records ready")

    docx_path = os.path.join(output_dir, "pipeline_brief.docx")
    print(f"  [3/4] Generating Word brief → {docx_path}")
    format_docx(summarised_deals, docx_path)
    print(f"        → {os.path.getsize(docx_path):,} bytes")

    pptx_path = os.path.join(output_dir, "pipeline_deck.pptx")
    print(f"  [4/4] Generating PowerPoint deck → {pptx_path}")
    format_pptx(summarised_deals, pptx_path)
    print(f"        → {os.path.getsize(pptx_path):,} bytes")

    print("── Done ────────────────────────────────────────────────\n")

    return {
        "slack":      slack_payload,
        "salesforce": sf_payloads,
        "docx_path":  docx_path,
        "pptx_path":  pptx_path,
    }
